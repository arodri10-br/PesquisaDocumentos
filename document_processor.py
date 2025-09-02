import os
from datetime import datetime
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
import PyPDF2
import pdfplumber
import pytesseract

try:
    from pdf2image import convert_from_path
except Exception:
    convert_from_path = None


class DocumentProcessor:
    SUPPORTED_EXTENSIONS = {
        ".doc": "doc",
        ".docx": "docx",
        ".ppt": "ppt",
        ".pptx": "pptx",
        ".xls": "xls",
        ".xlsx": "xlsx",
        ".pdf": "pdf",
        ".txt": "txt",
    }

    def __init__(self, poppler_path: str | None = None, tesseract_cmd: str | None = None):
        self.poppler_path = poppler_path
        if tesseract_cmd:
            pytesseract.pytesseract.tesseract_cmd = tesseract_cmd

    # -----------------------
    # Banco de dados (desacoplado)
    # -----------------------
    def scan_folder(self, folder_path: str, session, DocumentModel) -> int:
        count = 0
        folder_path = os.path.abspath(folder_path)

        for root, _, files in os.walk(folder_path):
            for fname in files:
                ext = os.path.splitext(fname)[1].lower()
                if ext not in self.SUPPORTED_EXTENSIONS:
                    continue

                file_path = os.path.abspath(os.path.join(root, fname))
                exists = session.query(DocumentModel).filter_by(filepath=file_path).first()
                if exists:
                    continue

                try:
                    st = os.stat(file_path)
                    created = datetime.fromtimestamp(st.st_ctime)
                    modified = datetime.fromtimestamp(st.st_mtime)
                    size = st.st_size
                except Exception:
                    created = modified = None
                    size = None

                doc = DocumentModel(
                    filename=fname,
                    filepath=file_path,
                    file_type=self.SUPPORTED_EXTENSIONS[ext],
                    file_size=size,
                    created_date=created,
                    modified_date=modified,
                    indexed_date=None,
                    status="pending",
                    folder_path=os.path.abspath(root),
                )
                session.add(doc)
                count += 1

        session.commit()
        return count

    # -----------------------
    # Extração de conteúdo
    # -----------------------
    def extract_content(self, file_path: str, file_type: str | None) -> str:
        try:
            ext = (file_type or os.path.splitext(file_path)[1].lstrip(".")).lower()

            if ext == "docx":
                return self._extract_docx(file_path)
            if ext == "doc":
                return "Conteúdo não disponível para arquivos .doc antigos"
            if ext in ("xlsx", "xls"):
                return self._extract_excel(file_path)
            if ext in ("pptx", "ppt"):
                return self._extract_powerpoint(file_path)
            if ext == "pdf":
                return self._extract_pdf(file_path)
            if ext == "txt":
                return self._extract_txt(file_path)
            return ""
        except Exception as e:
            print(f"[document_processor] Erro ao extrair {file_path}: {e}")
            return ""

    def _extract_txt(self, file_path: str) -> str:
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
        except Exception as e:
            print(f"[document_processor] TXT falhou {file_path}: {e}")
            return ""

    def _extract_docx(self, file_path: str) -> str:
        doc = DocxDocument(file_path)
        return "\n".join(p.text for p in doc.paragraphs if p.text)

    def _extract_excel(self, file_path: str) -> str:
        wb = load_workbook(file_path, data_only=True)
        parts: list[str] = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"--- Planilha: {sheet_name} ---")
            for row in ws.iter_rows(values_only=True):
                vals = [str(c) for c in row if c is not None]
                if vals:
                    parts.append("\t".join(vals))
        return "\n".join(parts)

    def _extract_powerpoint(self, file_path: str) -> str:
        prs = Presentation(file_path)
        parts: list[str] = []
        for i, slide in enumerate(prs.slides, start=1):
            parts.append(f"--- Slide {i} ---")
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False) and shape.text_frame:
                    txt = shape.text_frame.text
                    if txt:
                        parts.append(txt)
                elif hasattr(shape, "text"):
                    if shape.text:
                        parts.append(shape.text)
        return "\n".join(parts)

    def _extract_pdf(self, file_path: str) -> str:
        parts: list[str] = []

        # 1) pdfplumber (digital)
        try:
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    txt = page.extract_text(x_tolerance=2, y_tolerance=2) or ""
                    txt = txt.strip()
                    if txt:
                        parts.append(txt)
        except Exception as e:
            print(f"[document_processor] pdfplumber falhou {file_path}: {e}")

        # 2) PyPDF2 (fallback)
        if not parts:
            try:
                with open(file_path, "rb") as fh:
                    reader = PyPDF2.PdfReader(fh)
                    for p in reader.pages:
                        txt = p.extract_text() or ""
                        txt = txt.strip()
                        if txt:
                            parts.append(txt)
            except Exception as e:
                print(f"[document_processor] PyPDF2 falhou {file_path}: {e}")

        # 3) OCR (escaneado)
        if not parts and convert_from_path is not None:
            try:
                images = convert_from_path(file_path, dpi=300, poppler_path=self.poppler_path)
                ocr_parts = []
                for img in images:
                    txt = pytesseract.image_to_string(img, lang="por+eng", config="--oem 3 --psm 6")
                    if txt and txt.strip():
                        ocr_parts.append(txt.strip())
                if ocr_parts:
                    parts.extend(ocr_parts)
            except Exception as e:
                print(f"[document_processor] OCR falhou {file_path}: {e}")

        return "\n\n".join(parts).strip()
